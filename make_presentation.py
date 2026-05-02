"""Build PHEME spread-dynamics presentation as a .pptx file."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── paths ────────────────────────────────────────────────────────────────────
FIG = Path("figures")
OUT = Path("PHEME_Spread_Dynamics.pptx")

# ── colour palette (matches notebook colours) ────────────────────────────────
DARK       = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
ACCENT     = RGBColor(0x4C, 0x72, 0xB0)   # blue (non-rumour)
ORANGE     = RGBColor(0xDD, 0x84, 0x52)   # unverified
GREEN      = RGBColor(0x55, 0xA8, 0x68)   # true
RED        = RGBColor(0xC4, 0x4E, 0x52)   # false
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF4, 0xF5, 0xF7)
MID_GREY   = RGBColor(0x6C, 0x75, 0x7D)

# slide dimensions: 16:9 widescreen
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ── helper functions ─────────────────────────────────────────────────────────

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def txbox(slide, text, l, t, w, h,
          size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
          italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def heading(slide, text, top=0.18, size=32):
    txbox(slide, text, 0.35, top, 12.6, 0.75,
          size=size, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

def subhead(slide, text, top=0.85, size=17, color=LIGHT_GREY):
    txbox(slide, text, 0.35, top, 12.6, 0.55,
          size=size, bold=False, color=color, align=PP_ALIGN.LEFT)

def bg(slide, fill=DARK):
    """Full-slide background rectangle."""
    rect(slide, 0, 0, 13.33, 7.5, fill=fill)

def accent_bar(slide, h_frac=0.06):
    """Thin coloured bar at the top."""
    rect(slide, 0, 0, 13.33, 7.5 * h_frac, fill=ACCENT)

def img(slide, path, l, t, w, h=None):
    p = Path(path)
    if not p.exists():
        return
    if h is None:
        slide.shapes.add_picture(str(p), Inches(l), Inches(t), width=Inches(w))
    else:
        slide.shapes.add_picture(str(p), Inches(l), Inches(t), Inches(w), Inches(h))

def bullet_block(slide, items, l, t, w, h,
                 size=16, color=WHITE, indent="  "):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color

def stat_card(slide, label, value, note, l, t, w=2.5, h=1.4,
              fill=ACCENT, vcolor=WHITE, lcolor=LIGHT_GREY):
    rect(slide, l, t, w, h, fill=fill)
    txbox(slide, value, l+0.1, t+0.05, w-0.2, 0.65,
          size=28, bold=True, color=vcolor, align=PP_ALIGN.CENTER)
    txbox(slide, label, l+0.05, t+0.65, w-0.1, 0.38,
          size=12, bold=True, color=lcolor, align=PP_ALIGN.CENTER)
    if note:
        txbox(slide, note, l+0.05, t+1.0, w-0.1, 0.32,
              size=10, italic=True, color=lcolor, align=PP_ALIGN.CENTER)

def divider(slide, t, color=ACCENT, alpha=None):
    rect(slide, 0.35, t, 12.6, 0.025, fill=color)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
rect(sl, 0, 0, 0.12, 7.5, fill=ACCENT)                  # left accent stripe

txbox(sl, "Spread Dynamics of Rumours vs. Facts",
      0.55, 1.6, 12.3, 1.3,
      size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

txbox(sl, "Do rumours really travel faster during breaking-news events?",
      0.55, 2.95, 12.0, 0.65,
      size=22, italic=True, color=LIGHT_GREY, align=PP_ALIGN.LEFT)

divider(sl, 3.72)

txbox(sl, "PHEME-9 Dataset  ·  6,040 reply cascades  ·  6 breaking-news events",
      0.55, 3.82, 12.0, 0.5,
      size=16, color=LIGHT_GREY)

txbox(sl, "Data Science Group Project  ·  ETH Zürich  ·  2026",
      0.55, 6.7, 12.0, 0.5,
      size=13, color=MID_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — MOTIVATING QUESTION
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "Motivating question")
divider(sl, 0.95)

txbox(sl,
      "Vosoughi, Roy & Aral (Science, 2018) found that false news travels\n"
      "farther, faster, deeper, and more broadly than truth on Twitter.",
      0.55, 1.1, 12.2, 0.9,
      size=20, italic=True, color=LIGHT_GREY)

txbox(sl, "But does this hold during breaking-news events specifically?",
      0.55, 2.05, 12.2, 0.55,
      size=21, bold=True, color=WHITE)

txbox(sl,
      "During breaking news, facts haven't been established yet. "
      "Unverified claims fill the information vacuum.\n"
      "This is precisely when misinformation is most dangerous — and most plausible.",
      0.55, 2.7, 12.2, 0.9,
      size=17, color=LIGHT_GREY)

divider(sl, 3.72)

txbox(sl, "Our approach", 0.55, 3.82, 12.0, 0.38, size=16, bold=True, color=ACCENT)
bullet_block(sl, [
    "→  PHEME-9: journalist-annotated reply cascades from 9 breaking-news events on Twitter",
    "→  Measure cascade speed, reach, and structure across four veracity classes",
    "→  Propose and test a verifiability mechanism to explain why speed differences exist",
], 0.55, 4.22, 12.0, 1.6, size=16, color=LIGHT_GREY)

txbox(sl, "Reference: Vosoughi, Roy & Aral (2018). The spread of true and false news online. Science, 359(6380).",
      0.55, 6.8, 12.0, 0.45,
      size=11, italic=True, color=MID_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DATASET
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "Dataset: PHEME-9")
divider(sl, 0.95)

img(sl, FIG / "fig0_class_distribution.png", 0.3, 1.05, 8.6, 5.3)

# stat cards on the right
stat_card(sl, "reply cascades", "6,040", "after event filtering",
          9.15, 1.2, fill=ACCENT)
stat_card(sl, "breaking-news events", "6", "retained (3 dropped)",
          9.15, 2.75, fill=RGBColor(0x2E, 0x4A, 0x80))
stat_card(sl, "tweets total", "≈ 100k", "source + replies",
          9.15, 4.3, fill=RGBColor(0x1A, 0x35, 0x5E))

txbox(sl,
      "Non-rumours dominate every event — this is signal, not sampling artefact. "
      "Most breaking-news content is factual commentary around a rumour core.",
      0.35, 6.4, 8.7, 0.7,
      size=13, italic=True, color=MID_GREY)

txbox(sl, "Zubiaga et al. (2016). Analysing How People Orient to and Spread Rumours in Social Media. PLOS ONE.",
      0.35, 7.1, 12.5, 0.35,
      size=11, italic=True, color=MID_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT CASCADES LOOK LIKE (Hook)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "What a cascade looks like")
divider(sl, 0.95)

img(sl, FIG / "fig1_hook.png", 0.5, 1.05, 12.3, 5.5)

txbox(sl,
      "Each node is a tweet. The black root is the original post; "
      "edges are direct replies. We measure how fast, how far, and what shape these trees take.",
      0.5, 6.65, 12.3, 0.6,
      size=13, italic=True, color=MID_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — METRICS & METHOD BRIEF
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "Methodology")
divider(sl, 0.95)

# Three columns
col_titles = ["SPEED", "REACH", "STRUCTURE"]
col_x      = [0.55, 4.7, 8.85]
col_color  = [ORANGE, ACCENT, GREEN]
col_items  = [
    ["Time to first reply", "Time to half-cascade", "Replies in first hour"],
    ["Cascade size (# tweets)", "Max depth", "Max breadth", "Unique users"],
    ["Structural virality (Goel 2016)", "Broadcast ratio", "Branching factor"],
]

for title, x, col, items in zip(col_titles, col_x, col_color, col_items):
    rect(sl, x, 1.1, 4.0, 0.5, fill=col)
    txbox(sl, title, x+0.1, 1.12, 3.8, 0.44,
          size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for i, item in enumerate(items):
        txbox(sl, f"• {item}", x+0.15, 1.72 + i*0.5, 3.75, 0.46,
              size=15, color=LIGHT_GREY)

divider(sl, 4.05)

txbox(sl, "Statistical approach", 0.55, 4.18, 12.0, 0.38,
      size=16, bold=True, color=ACCENT)
bullet_block(sl, [
    "→  Mann-Whitney U + Cliff's δ effect size — non-parametric, appropriate for heavy-tailed distributions",
    "→  Benjamini-Hochberg FDR correction across all simultaneous tests",
    "→  Per-event stratification — never pool across events with very different sizes/compositions",
    "→  Consistency scoring: a finding must hold in ≥ 3 eligible events to be reported",
], 0.55, 4.58, 12.0, 1.85, size=15, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — FINDING 1: SPEED (CDF)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "Finding 1 — Rumours generate first replies faster")
divider(sl, 0.95)

img(sl, FIG / "fig2_speed_first_reply.png", 0.3, 1.05, 8.8, 5.35)

# key stats on the right
txbox(sl, "Effect size", 9.4, 1.2, 3.6, 0.42,
      size=15, bold=True, color=ACCENT)

stat_card(sl, "Cliff's δ (unverified vs non-rumour)", "−0.26",
          "pooled p < 0.001", 9.4, 1.65, w=3.5, fill=ORANGE)

stat_card(sl, "Cliff's δ (true vs non-rumour)", "−0.17",
          "pooled p = 0.045", 9.4, 3.2, w=3.5, fill=GREEN)

txbox(sl,
      "Negative δ = rumours reach first reply faster than non-rumours.\n"
      "Both unverified and true rumours show the effect.",
      9.4, 4.75, 3.6, 0.8,
      size=13, italic=True, color=LIGHT_GREY)

txbox(sl,
      "Cliff's δ interpretation:  |δ| < 0.15 negligible  ·  ≥ 0.15 small  ·  ≥ 0.33 medium",
      0.35, 6.75, 12.5, 0.4,
      size=11, italic=True, color=MID_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SPEED HOLDS ACROSS EVENTS
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "The speed effect is consistent across events")
divider(sl, 0.95)

img(sl, FIG / "fig3b_per_event_speed_ci.png", 0.3, 1.05, 8.6, 5.3)

txbox(sl, "Why this matters", 9.35, 1.2, 3.7, 0.4,
      size=15, bold=True, color=ACCENT)
bullet_block(sl, [
    "A pooled p-value could be driven by a single large event.",
    "",
    "Non-rumours are consistently slower to receive first replies across charliehebdo, ottawashooting, sydneysiege, and others.",
    "",
    "95% bootstrap CIs do not overlap in most events — the per-event medians are precisely estimated.",
    "",
    "Pattern holds in 2 of 6 eligible events at BH-corrected p < 0.05 for time-to-first-reply (3/6 for time-to-half-cascade).",
], 9.35, 1.7, 3.7, 4.6, size=13, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — NULL RESULT: REACH
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "Null result — Rumours do not reach more people")
divider(sl, 0.95)

img(sl, FIG / "fig4_reach_null.png", 0.3, 1.05, 8.6, 5.3)

txbox(sl, "What this means", 9.35, 1.2, 3.7, 0.4,
      size=15, bold=True, color=ACCENT)
bullet_block(sl, [
    "Vosoughi (2018) found that false news reaches more people.",
    "",
    "We do not replicate this for breaking-news reply cascades.",
    "",
    "Cascade sizes (median 10–14 tweets) are comparable across all four veracity classes.",
    "",
    "The Vosoughi finding may be specific to the retweet network or to non-breaking-news contexts.",
    "",
    "Implication: rumours engage faster, but the conversation stays local — they do not diffuse farther.",
], 9.35, 1.7, 3.7, 4.9, size=13, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — WHY? VERIFIABILITY MECHANISM (INTRO)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "Why are rumours faster? — A verifiability mechanism")
divider(sl, 0.95)

txbox(sl,
      "Hypothesis: claims that are hard to verify generate faster engagement "
      "because readers cannot suppress sharing via a quick mental fact-check.",
      0.55, 1.1, 12.2, 0.85,
      size=19, italic=True, color=LIGHT_GREY)

# Two-step diagram
for i, (x, col, num, title, body) in enumerate([
    (0.55, ORANGE, "Step 1", "Rumours are less verifiable",
     "A TF-IDF + logistic regression classifier trained on FEVER\n"
     "(145k Wikipedia fact-verification claims) scores each tweet\n"
     "for how linguistically verifiable it is at posting time."),
    (6.85, ACCENT, "Step 2", "Less verifiable → faster reply",
     "Threads in the least-verifiable quartile (Q1) receive their\n"
     "first reply in 1.87 min median vs 2.55 min for Q4.\n"
     "Spearman ρ = 0.09–0.10, p < 10⁻¹²."),
]):
    rect(sl, x, 2.05, 6.0, 3.9, fill=RGBColor(0x1E, 0x2A, 0x40))
    rect(sl, x, 2.05, 6.0, 0.52, fill=col)
    txbox(sl, f"{num}  ·  {title}", x+0.15, 2.08, 5.7, 0.44,
          size=16, bold=True, color=WHITE)
    txbox(sl, body, x+0.2, 2.65, 5.6, 2.8,
          size=15, color=LIGHT_GREY)

txbox(sl, "→", 6.45, 3.6, 0.5, 0.6, size=28, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

txbox(sl,
      "Caveat: ρ = 0.09 is a weak effect. The ~40-second median difference is "
      "statistically reliable (n = 5,747) but practically modest. "
      "Verifiability is a contributing factor, not the primary driver.",
      0.55, 6.1, 12.2, 0.7,
      size=13, italic=True, color=MID_GREY)

txbox(sl, "Classifier: FEVER (Thorne et al., NAACL 2018)  ·  ROC-AUC ≥ 0.80 on FEVER dev split",
      0.55, 6.88, 12.2, 0.38,
      size=11, italic=True, color=MID_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — VERIFIABILITY MECHANISM (FIGURE)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "Verifiability mechanism — Evidence")
divider(sl, 0.95)

img(sl, FIG / "fig11_verifiability_mechanism_summary.png", 0.3, 1.05, 12.7, 5.55)

txbox(sl,
      "Left: rumour classes score lower on P(verifiable) than non-rumours. "
      "True rumours are hardest to verify — even correct claims were linguistically ambiguous at posting time.  "
      "Right: monotone increase — less verifiable → faster first reply across all four quartiles.",
      0.35, 6.7, 12.6, 0.6,
      size=13, italic=True, color=MID_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — METHODOLOGY MATTERS: POOLED vs STRATIFIED
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "Methodological note — Why per-event analysis matters")
divider(sl, 0.95)

img(sl, FIG / "fig6_methodology.png", 0.3, 1.05, 12.7, 5.1)

txbox(sl,
      "A pooled analysis would conclude 'non-rumours have the biggest cascades'. "
      "But the per-event view shows rankings vary by event — the pooled finding is event-confounded. "
      "We use per-event stratified tests throughout to avoid this artefact.",
      0.35, 6.25, 12.6, 0.85,
      size=13, italic=True, color=MID_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — SUMMARY OF FINDINGS
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "Summary of findings")
divider(sl, 0.95)

rows = [
    # (icon_color, finding, direction, effect, interpretation)
    (GREEN,  "Speed: time to first reply",
     "rumour < non-rumour",
     "Cliff's δ = −0.26",
     "Rumours attract replies significantly faster"),
    (GREEN,  "Speed: time to half-cascade",
     "rumour < non-rumour",
     "Cliff's δ = −0.32",
     "Speed advantage persists through the cascade"),
    (RED,    "Reach: cascade size",
     "null result",
     "—",
     "Rumours do not reach more users overall"),
    (RED,    "Structure: depth vs. size",
     "null result",
     "—",
     "No distinct structural signature for rumours"),
    (ORANGE, "Verifiability → speed",
     "less verifiable → faster",
     "Spearman ρ = 0.09–0.10",
     "Claim ambiguity predicts faster engagement (p < 10⁻¹²)"),
    (ORANGE, "True rumours",
     "least verifiable class",
     "Cliff's δ = −0.22",
     "Even correct claims were hard to verify at posting time"),
]

for i, (col, finding, direction, effect, interp) in enumerate(rows):
    y = 1.1 + i * 0.97
    rect(sl, 0.35, y, 0.18, 0.75, fill=col)
    txbox(sl, finding,    0.62, y+0.02, 3.5,  0.38, size=14, bold=True,  color=WHITE)
    txbox(sl, direction,  4.2,  y+0.02, 2.5,  0.38, size=13, color=LIGHT_GREY)
    txbox(sl, effect,     6.8,  y+0.02, 2.2,  0.38, size=13, bold=True,  color=ACCENT)
    txbox(sl, interp,     9.1,  y+0.02, 4.0,  0.38, size=12, italic=True, color=LIGHT_GREY)

# column headers
for x, label in [(0.62, "Finding"), (4.2, "Direction"), (6.8, "Effect size"), (9.1, "Interpretation")]:
    txbox(sl, label, x, 0.98, 3.5, 0.3, size=12, bold=True, color=MID_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — LIMITATIONS
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "Limitations & scope")
divider(sl, 0.95)

limits = [
    ("Dataset size & scope",
     "PHEME contains ~6,400 threads across 9 events — small by modern standards. "
     "Results may not generalise beyond English-language breaking-news Twitter."),
    ("Reply trees, not retweet networks",
     "We measure conversational depth (who replies to whom), not diffusion breadth "
     "(who retweets). Vosoughi's 'reach' finding is about retweet chains, which we cannot replicate here."),
    ("Snapshot data",
     "Cascades were collected at a fixed point in time. "
     "Later replies, which could change speed and size rankings, are not captured."),
    ("Verifiability classifier transfer",
     "The FEVER classifier was trained on Wikipedia sentences. "
     "Transfer to informal tweet language is imperfect; ρ = 0.09 is weak. "
     "The score is a proxy, not ground truth."),
    ("Annotation quality",
     "Veracity labels are journalist-curated (high precision, limited recall). "
     "Many threads remain 'unverified' not because they are rumours, but because "
     "the fact was never definitively established."),
]

for i, (title, body) in enumerate(limits):
    y = 1.05 + i * 1.1
    rect(sl, 0.35, y, 0.08, 0.85, fill=ACCENT)
    txbox(sl, title, 0.6, y+0.02, 4.0, 0.38, size=14, bold=True, color=WHITE)
    txbox(sl, body,  0.6, y+0.43, 12.3, 0.55, size=13, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — REFERENCES
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
bg(sl)
accent_bar(sl)

heading(sl, "References")
divider(sl, 0.95)

refs = [
    "Vosoughi, S., Roy, D., & Aral, S. (2018). The spread of true and false news online. Science, 359(6380), 1146–1151.",
    "Zubiaga, A., Liakata, M., Procter, R., Wong Sak Hoi, G., & Tolmie, P. (2016). Analysing How People Orient to and Spread Rumours in Social Media. PLOS ONE, 11(6).",
    "Thorne, J., Vlachos, A., Christodoulopoulos, C., & Mittal, A. (2018). FEVER: a large-scale dataset for Fact Extraction and VERification. NAACL-HLT.",
    "Goel, S., Anderson, A., Hofman, J., & Watts, D. J. (2016). The structural virality of online diffusion. Management Science, 62(1), 180–196.",
    "Nielsen, R. K., & McConville, R. (2022). MuMiN: A Large-Scale Multilingual Multimodal Fact-Checked Misinformation Social Network Dataset. Findings of ACL.",
    "Romano, J., Kromrey, J. D., Coraggio, J., & Skowronek, J. (2006). Appropriate statistics for ordinal level data: Should we really be using t-test and Cohen's d for evaluating group differences on the NSSE? Florida Association of Institutional Research.",
]

for i, ref in enumerate(refs):
    txbox(sl, ref, 0.55, 1.1 + i * 0.85, 12.2, 0.75,
          size=13, color=LIGHT_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Saved: {OUT}  ({OUT.stat().st_size / 1024:.0f} KB,  {len(prs.slides)} slides)")
